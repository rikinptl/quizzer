#!/usr/bin/env python3
"""
Text extraction script for PDF and PowerPoint files.
Supports PDF, PPT, and PPTX formats.
"""

import sys
import os
import PyPDF2
from pptx import Presentation
import argparse
import logging

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def extract_pdf_text(file_path):
    """Extract text from PDF file."""
    try:
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n"
                
        logger.info(f"Successfully extracted text from PDF: {len(text)} characters")
        return text.strip()
        
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path}: {str(e)}")
        raise

def extract_pptx_text(file_path):
    """Extract text from PowerPoint file (PPTX)."""
    try:
        text = ""
        presentation = Presentation(file_path)
        
        for slide_num, slide in enumerate(presentation.slides):
            text += f"\n--- Slide {slide_num + 1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
                    
        logger.info(f"Successfully extracted text from PPTX: {len(text)} characters")
        return text.strip()
        
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {str(e)}")
        raise

def extract_text(file_path):
    """Extract text from supported file formats."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")
    
    file_extension = os.path.splitext(file_path)[1].lower()
    
    if file_extension == '.pdf':
        return extract_pdf_text(file_path)
    elif file_extension in ['.ppt', '.pptx']:
        return extract_pptx_text(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}. Supported formats: .pdf, .ppt, .pptx")

def clean_text(text):
    """Clean and format extracted text."""
    if not text:
        return ""
    
    # Remove excessive whitespace
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        cleaned_line = line.strip()
        if cleaned_line:  # Skip empty lines
            cleaned_lines.append(cleaned_line)
    
    # Join lines with single newlines
    cleaned_text = '\n'.join(cleaned_lines)
    
    # Remove excessive spaces
    import re
    cleaned_text = re.sub(r' +', ' ', cleaned_text)
    
    return cleaned_text

def main():
    parser = argparse.ArgumentParser(description='Extract text from PDF or PowerPoint files')
    parser.add_argument('file_path', help='Path to the PDF or PowerPoint file')
    parser.add_argument('--output', '-o', help='Output file path (default: stdout)')
    parser.add_argument('--clean', '-c', action='store_true', help='Clean and format the extracted text')
    
    args = parser.parse_args()
    
    try:
        # Extract text
        text = extract_text(args.file_path)
        
        # Clean text if requested
        if args.clean:
            text = clean_text(text)
        
        # Output text
        if args.output:
            with open(args.output, 'w', encoding='utf-8') as f:
                f.write(text)
            logger.info(f"Text saved to: {args.output}")
        else:
            print(text)
            
    except Exception as e:
        logger.error(f"Failed to extract text: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()
